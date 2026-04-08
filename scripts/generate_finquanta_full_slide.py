# -*- coding: utf-8 -*-
"""Generate a one-slide PPTX encompassing AI, Classic Strategies, and Quantum Optimization."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(0x02, 0x05, 0x12)
PANEL = RGBColor(0x05, 0x08, 0x16)
PANEL_BORDER = RGBColor(0x2D, 0x6D, 0xF0)
BLUE = RGBColor(0x4F, 0x79, 0xFF)
CYAN = RGBColor(0x00, 0xC8, 0xFF)
TEXT = RGBColor(0xE8, 0xF1, 0xFF)
MUTED = RGBColor(0xA8, 0xBC, 0xD6)
EMPH = RGBColor(0x00, 0xF0, 0xFF)


def _send_to_back(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(2, el)


def _textbox(slide, x, y, w, h, text, *, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def _panel(slide, x, y, w, h, title):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = PANEL_BORDER
    panel.line.width = Pt(1.6)

    title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y + Inches(0.04), w - Inches(0.1), Inches(0.42))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(0x0B, 0x12, 0x2A)
    title_box.line.color.rgb = PANEL_BORDER
    title_box.line.width = Pt(1.0)
    _textbox(slide, x + Inches(0.08), y + Inches(0.1), w - Inches(0.16), Inches(0.24), title, size=15, bold=True, align=PP_ALIGN.CENTER)
    return panel


def _chip(slide, x, y, w, h, text, *, fill=RGBColor(0xE9, 0xEE, 0xFF), text_color=RGBColor(0x12, 0x1C, 0x30), size=10, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PANEL_BORDER
    shape.line.width = Pt(0.8)
    _textbox(slide, x + Inches(0.02), y + Inches(0.03), w - Inches(0.04), h - Inches(0.06), text, size=size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return shape


def _bottom_arrow(slide, x, y, w, h, text):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x0B, 0xB2, 0xFF)
    arrow.line.color.rgb = RGBColor(0x0B, 0xB2, 0xFF)
    _textbox(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), h - Inches(0.1), text, size=12, bold=True, align=PP_ALIGN.CENTER)


def _connector(slide, x1, y1, x2, y2, color=RGBColor(0x4F, 0x79, 0xFF), width=Pt(1.2)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    return line


def main():
    root = Path(__file__).resolve().parents[1]
    out = root / "doc" / "FinQuanta_AI策略量子融合_一页.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else -1])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(prs.slide_width), int(prs.slide_height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    _send_to_back(slide, bg)

    # Subtle glow blocks
    for x, y, w, h, color in [
        (Inches(-0.2), Inches(0.0), Inches(1.0), Inches(7.5), RGBColor(0x0B, 0x33, 0xA8)),
        (Inches(12.5), Inches(0.0), Inches(1.0), Inches(7.5), RGBColor(0x0B, 0x33, 0xA8)),
        (Inches(10.8), Inches(0.2), Inches(2.1), Inches(2.0), RGBColor(0x13, 0x55, 0xFF)),
    ]:
        glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        glow.fill.solid()
        glow.fill.fore_color.rgb = color
        glow.fill.transparency = 0.78
        glow.line.fill.background()
        _send_to_back(slide, glow)

    _textbox(slide, Inches(0.38), Inches(0.1), Inches(12.6), Inches(0.32), "FinQuanta：AI 智能决策 × 经典选股策略 × 量子组合优化", size=16, bold=True, color=RGBColor(0xB7, 0xD7, 0xFF), align=PP_ALIGN.CENTER)

    left_x, left_y, left_w, left_h = Inches(0.25), Inches(0.55), Inches(3.45), Inches(5.45)
    center_x, center_y, center_w, center_h = Inches(4.0), Inches(0.55), Inches(5.2), Inches(5.45)
    right_x, right_y, right_w, right_h = Inches(9.55), Inches(0.55), Inches(3.25), Inches(5.45)

    # ------------------ LEFT PANEL: CORE CHALLENGES ------------------
    _panel(slide, left_x, left_y, left_w, left_h, "智能金融与量化交易挑战")
    _textbox(
        slide,
        left_x + Inches(0.2),
        left_y + Inches(0.62),
        left_w - Inches(0.4),
        Inches(0.7),
        "现代金融市场具有高噪音、非线性和多重博弈特征，使得传统纯规则或纯人工决策面临巨大的性能与广度瓶颈。",
        size=10,
        color=MUTED,
    )
    alert = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x + Inches(0.18), left_y + Inches(1.3), left_w - Inches(0.36), Inches(2.5))
    alert.fill.solid()
    alert.fill.fore_color.rgb = RGBColor(0x14, 0x7E, 0xC1)
    alert.line.color.rgb = RGBColor(0x59, 0xD1, 0xFF)
    
    _textbox(
        slide,
        left_x + Inches(0.25),
        left_y + Inches(1.4),
        left_w - Inches(0.5),
        Inches(0.4),
        "核心痛点：信息过载、策略固化、最优解难求",
        size=10,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    left_bullets = [
        "信息与决策不对等：全市场标的与新闻海量，传统系统缺乏自主学习与意图识别能力。",
        "策略适应性弱：固定因子的经典选股策略在剧烈震荡市中易失效，需要AI动态择时。",
        "组合规模化瓶颈：当候选池扩张时(如500选30)，超10^48规模的解空间使经典算力陷入局部最优。",
    ]
    for idx, item in enumerate(left_bullets):
        _textbox(
            slide,
            left_x + Inches(0.25),
            left_y + Inches(1.85 + 0.6 * idx),
            left_w - Inches(0.5),
            Inches(0.55),
            f"• {item}",
            size=9,
            color=TEXT,
        )

    _textbox(
        slide,
        left_x + Inches(0.2),
        left_y + Inches(4.0),
        left_w - Inches(0.4),
        Inches(1.2),
        "全新范式：\n利用 AI 大模型实现自主意图识别与宏观择时；\n通过多维经典策略生成高质量Alpha候选标的；\n引入量子计算 QAOA 打破维度诅咒完成全局寻优。",
        size=10,
        bold=True,
        color=RGBColor(0x70, 0xE8, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    # ------------------ CENTER PANEL: THREE-LAYER ARCHITECTURE ------------------
    _panel(slide, center_x, center_y, center_w, center_h, "AI × 策略 × 量子 一体化架构")
    
    # Layer 1: AI
    y_ai = center_y + Inches(0.65)
    _chip(slide, center_x + Inches(0.2), y_ai, center_w - Inches(0.4), Inches(1.15), "", fill=RGBColor(0x09, 0x18, 0x3D))
    _textbox(slide, center_x + Inches(0.2), y_ai + Inches(0.05), center_w - Inches(0.4), Inches(0.2), "【顶层】AI 智能决策与意图网络 (OpenClaw Agent)", size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    
    _chip(slide, center_x + Inches(0.35), y_ai + Inches(0.45), Inches(1.3), Inches(0.5), "自然语言意图交互", fill=RGBColor(0x11, 0x2A, 0x66), text_color=TEXT, size=9)
    _chip(slide, center_x + Inches(1.95), y_ai + Inches(0.45), Inches(1.3), Inches(0.5), "宏观情绪与择时", fill=RGBColor(0x11, 0x2A, 0x66), text_color=TEXT, size=9)
    _chip(slide, center_x + Inches(3.55), y_ai + Inches(0.45), Inches(1.3), Inches(0.5), "跨Agent协作网络", fill=RGBColor(0x11, 0x2A, 0x66), text_color=TEXT, size=9)

    _connector(slide, center_x + center_w/2, y_ai + Inches(1.15), center_x + center_w/2, y_ai + Inches(1.4), width=Pt(2.0), color=CYAN)

    # Layer 2: Classic Strategies
    y_str = center_y + Inches(2.05)
    _chip(slide, center_x + Inches(0.2), y_str, center_w - Inches(0.4), Inches(1.15), "", fill=RGBColor(0x0E, 0x26, 0x4D))
    _textbox(slide, center_x + Inches(0.2), y_str + Inches(0.05), center_w - Inches(0.4), Inches(0.2), "【中层】多维经典选股策略矩阵", size=11, bold=True, color=RGBColor(0xFF, 0xCC, 0x66), align=PP_ALIGN.CENTER)
    
    _chip(slide, center_x + Inches(0.35), y_str + Inches(0.45), Inches(1.3), Inches(0.5), "多因子Alpha挖掘", fill=RGBColor(0x26, 0x3F, 0x73), text_color=TEXT, size=9)
    _chip(slide, center_x + Inches(1.95), y_str + Inches(0.45), Inches(1.3), Inches(0.5), "形态与事件驱动", fill=RGBColor(0x26, 0x3F, 0x73), text_color=TEXT, size=9)
    _chip(slide, center_x + Inches(3.55), y_str + Inches(0.45), Inches(1.3), Inches(0.5), "协方差与特征提取", fill=RGBColor(0x26, 0x3F, 0x73), text_color=TEXT, size=9)

    _connector(slide, center_x + center_w/2, y_str + Inches(1.15), center_x + center_w/2, y_str + Inches(1.4), width=Pt(2.0), color=RGBColor(0xFF, 0xCC, 0x66))

    # Layer 3: Quantum
    y_q = center_y + Inches(3.45)
    _chip(slide, center_x + Inches(0.2), y_q, center_w - Inches(0.4), Inches(1.15), "", fill=RGBColor(0x1B, 0x1A, 0x4A))
    _textbox(slide, center_x + Inches(0.2), y_q + Inches(0.05), center_w - Inches(0.4), Inches(0.2), "【底层】量子计算组合优化核心", size=11, bold=True, color=RGBColor(0xFF, 0x99, 0xFF), align=PP_ALIGN.CENTER)
    
    _chip(slide, center_x + Inches(0.35), y_q + Inches(0.45), Inches(1.3), Inches(0.5), "QUBO 二次建模", fill=RGBColor(0x3B, 0x2A, 0x6A), text_color=TEXT, size=9)
    _chip(slide, center_x + Inches(1.95), y_q + Inches(0.45), Inches(1.3), Inches(0.5), "QAOA 并行寻优", fill=RGBColor(0x6B, 0x3A, 0xAA), text_color=TEXT, size=9, bold=True)
    _chip(slide, center_x + Inches(3.55), y_q + Inches(0.45), Inches(1.3), Inches(0.5), "持仓/风险硬约束", fill=RGBColor(0x3B, 0x2A, 0x6A), text_color=TEXT, size=9)

    # Bottom summary
    flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_x + Inches(0.22), center_y + Inches(4.75), center_w - Inches(0.44), Inches(0.55))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0x09, 0x0F, 0x24)
    flow_box.line.color.rgb = PANEL_BORDER
    flow_box.line.width = Pt(1.0)
    _textbox(
        slide,
        center_x + Inches(0.25),
        center_y + Inches(4.85),
        center_w - Inches(0.5),
        Inches(0.4),
        "全自动化投研闭环：大模型指引方向 → 经典策略圈定标的池 → 量子优化生成最优配置",
        size=10,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    # ------------------ RIGHT PANEL: KEY RESULTS ------------------
    _panel(slide, right_x, right_y, right_w, right_h, "多维融合关键优势")
    
    stat_titles = [("24 / 7", "AI全天候动态自适应"), ("多因子", "经典与前沿特征同源"), ("1.20 夏普", "QAOA高维解质量提升")]
    stat_xs = [right_x + Inches(0.1), right_x + Inches(1.15), right_x + Inches(2.2)]
    
    for i, (big, small) in enumerate(stat_titles):
        _textbox(slide, right_x + Inches(0.1), right_y + Inches(0.8) + i * Inches(1.3), Inches(1.2), Inches(0.5), big, size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        _textbox(slide, right_x + Inches(1.4), right_y + Inches(0.9) + i * Inches(1.3), Inches(1.7), Inches(0.5), small, size=11, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
        if i == 0:
            _textbox(slide, right_x + Inches(1.4), right_y + Inches(1.15) + i * Inches(1.3), Inches(1.7), Inches(0.6), "Agent根据实时盘面和新闻自动干预仓位。", size=8, color=MUTED)
        elif i == 1:
            _textbox(slide, right_x + Inches(1.4), right_y + Inches(1.15) + i * Inches(1.3), Inches(1.7), Inches(0.6), "支持微观均值回归、动量破局等数十种算子。", size=8, color=MUTED)
        elif i == 2:
            _textbox(slide, right_x + Inches(1.4), right_y + Inches(1.15) + i * Inches(1.3), Inches(1.7), Inches(0.6), "仿真中量子求解相对经典贪心提升夏普比2.6%。", size=8, color=MUTED)

        if i < 2:
            _connector(slide, right_x + Inches(0.2), right_y + Inches(1.9) + i * Inches(1.3), right_x + right_w - Inches(0.2), right_y + Inches(1.9) + i * Inches(1.3), color=RGBColor(0x1F, 0x33, 0x66), width=Pt(1.0))

    _textbox(
        slide,
        right_x + Inches(0.2),
        right_y + Inches(4.55),
        right_w - Inches(0.4),
        Inches(0.52),
        "通过 “AI认知 + 策略生成 + 量子优化” 的三维驱动，真正让金融量化从“人工设定规则”迈向“机器自主进化”。",
        size=10,
        bold=True,
        color=RGBColor(0xFF, 0xCC, 0x66),
        align=PP_ALIGN.CENTER,
    )

    # ------------------ BOTTOM LABELS ------------------
    _bottom_arrow(slide, Inches(0.45), Inches(6.12), Inches(3.2), Inches(0.44), "金融痛点：复杂市场与算力瓶颈")
    _bottom_arrow(slide, Inches(4.45), Inches(6.12), Inches(4.35), Inches(0.44), "FinQuanta 平台：三维一体全栈架构")
    _bottom_arrow(slide, Inches(9.7), Inches(6.12), Inches(2.9), Inches(0.44), "量化投资的智能进化新范式")

    _textbox(
        slide,
        Inches(0.35),
        Inches(6.68),
        Inches(12.65),
        Inches(0.25),
        "说明：架构与数据展示基于 FinQuanta 平台设计与端到端测试，具体投资表现视实际运行策略为准。",
        size=8,
        color=RGBColor(0x8C, 0x9B, 0xAF),
        align=PP_ALIGN.CENTER,
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print("Wrote", out)


if __name__ == "__main__":
    main()
