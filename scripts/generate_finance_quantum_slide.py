# -*- coding: utf-8 -*-
"""Generate a one-slide PPTX in a three-panel launch-event style."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(0x02, 0x05, 0x12)
PANEL = RGBColor(0x05, 0x08, 0x16)
PANEL_BORDER = RGBColor(0x2D, 0x6D, 0xF0)
BLUE = RGBColor(0x4F, 0x79, 0xFF)
CYAN = RGBColor(0x00, 0xC8, 0xFF)
TEXT = RGBColor(0xE8, 0xF1, 0xFF)
MUTED = RGBColor(0xA8, 0xBC, 0xD6)
EMPH = RGBColor(0x00, 0xF0, 0xFF)


def _send_to_back(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(2, el)


def _textbox(slide, x, y, w, h, text, *, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def _panel(slide, x, y, w, h, title):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = PANEL_BORDER
    panel.line.width = Pt(1.6)

    title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y + Inches(0.04), w - Inches(0.1), Inches(0.42))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(0x0B, 0x12, 0x2A)
    title_box.line.color.rgb = PANEL_BORDER
    title_box.line.width = Pt(1.0)
    _textbox(slide, x + Inches(0.08), y + Inches(0.1), w - Inches(0.16), Inches(0.24), title, size=15, bold=True, align=PP_ALIGN.CENTER)
    return panel


def _chip(slide, x, y, w, h, text, *, fill=RGBColor(0xE9, 0xEE, 0xFF), text_color=RGBColor(0x12, 0x1C, 0x30), size=10, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PANEL_BORDER
    shape.line.width = Pt(0.8)
    _textbox(slide, x + Inches(0.02), y + Inches(0.03), w - Inches(0.04), h - Inches(0.06), text, size=size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return shape


def _bottom_arrow(slide, x, y, w, h, text):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x0B, 0xB2, 0xFF)
    arrow.line.color.rgb = RGBColor(0x0B, 0xB2, 0xFF)
    _textbox(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), h - Inches(0.1), text, size=12, bold=True, align=PP_ALIGN.CENTER)


def _connector(slide, x1, y1, x2, y2, color=RGBColor(0x4F, 0x79, 0xFF), width=Pt(1.2)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    return line


def _draw_chart(slide, x, y, w, h):
    chart = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    chart.fill.solid()
    chart.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFB)
    chart.line.color.rgb = RGBColor(0xE1, 0xE7, 0xEF)

    left = x + Inches(0.18)
    bottom = y + h - Inches(0.18)
    right = x + w - Inches(0.15)
    top = y + Inches(0.18)

    _connector(slide, left, bottom, right, bottom, color=RGBColor(0xB4, 0xBF, 0xCF), width=Pt(1.0))
    _connector(slide, left, bottom, left, top, color=RGBColor(0xB4, 0xBF, 0xCF), width=Pt(1.0))

    grid_y = [0.75, 0.55, 0.35]
    for r in grid_y:
        gy = y + Inches(r)
        _connector(slide, left, gy, right, gy, color=RGBColor(0xE3, 0xE8, 0xF0), width=Pt(0.8))

    _textbox(slide, x + Inches(0.12), y + Inches(0.03), w - Inches(0.24), Inches(0.18), "Search Complexity Comparison", size=8, color=RGBColor(0x46, 0x4F, 0x64), align=PP_ALIGN.CENTER)
    _textbox(slide, x + Inches(0.04), y + Inches(0.35), Inches(0.14), Inches(0.18), "高", size=7, color=RGBColor(0x6C, 0x75, 0x88), align=PP_ALIGN.CENTER)
    _textbox(slide, x + Inches(0.04), bottom - Inches(0.1), Inches(0.14), Inches(0.18), "低", size=7, color=RGBColor(0x6C, 0x75, 0x88), align=PP_ALIGN.CENTER)
    _textbox(slide, right - Inches(0.4), bottom + Inches(0.02), Inches(0.36), Inches(0.16), "规模", size=7, color=RGBColor(0x6C, 0x75, 0x88), align=PP_ALIGN.RIGHT)

    blue_points = [
        (left + Inches(0.05), bottom - Inches(0.04)),
        (left + Inches(0.55), bottom - Inches(0.06)),
        (left + Inches(1.05), bottom - Inches(0.08)),
        (left + Inches(1.55), bottom - Inches(0.15)),
        (left + Inches(2.0), bottom - Inches(0.48)),
        (left + Inches(2.25), bottom - Inches(0.98)),
    ]
    orange_points = [
        (left + Inches(0.05), bottom - Inches(0.03)),
        (left + Inches(0.6), bottom - Inches(0.09)),
        (left + Inches(1.1), bottom - Inches(0.19)),
        (left + Inches(1.6), bottom - Inches(0.34)),
        (left + Inches(2.0), bottom - Inches(0.6)),
        (left + Inches(2.25), bottom - Inches(1.1)),
    ]
    for pts, color in ((orange_points, RGBColor(0xFF, 0xA5, 0x3A)), (blue_points, RGBColor(0x2D, 0x73, 0xFF))):
        for p1, p2 in zip(pts, pts[1:]):
            _connector(slide, p1[0], p1[1], p2[0], p2[1], color=color, width=Pt(2.0))

    _textbox(slide, x + Inches(0.48), y + Inches(0.22), Inches(0.8), Inches(0.15), "量电混合", size=7, color=RGBColor(0xFF, 0xA5, 0x3A))
    _textbox(slide, x + Inches(1.35), y + Inches(0.22), Inches(0.9), Inches(0.15), "经典穷举/搜索", size=7, color=RGBColor(0x2D, 0x73, 0xFF))


def main():
    root = Path(__file__).resolve().parents[1]
    out = root / "doc" / "FinQuanta_金融量子应用_一页_三栏优化版.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else -1])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(prs.slide_width), int(prs.slide_height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    _send_to_back(slide, bg)

    # Subtle glow blocks to make the slide closer to the reference style.
    for x, y, w, h, color in [
        (Inches(-0.2), Inches(0.0), Inches(1.0), Inches(7.5), RGBColor(0x0B, 0x33, 0xA8)),
        (Inches(12.5), Inches(0.0), Inches(1.0), Inches(7.5), RGBColor(0x0B, 0x33, 0xA8)),
        (Inches(10.8), Inches(0.2), Inches(2.1), Inches(2.0), RGBColor(0x13, 0x55, 0xFF)),
    ]:
        glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        glow.fill.solid()
        glow.fill.fore_color.rgb = color
        glow.fill.transparency = 0.78
        glow.line.fill.background()
        _send_to_back(slide, glow)

    _textbox(slide, Inches(0.38), Inches(0.1), Inches(12.6), Inches(0.32), "应用价值：赋能产业创新发展｜金融场景补充", size=14, bold=True, color=RGBColor(0xB7, 0xD7, 0xFF), align=PP_ALIGN.CENTER)

    left_x, left_y, left_w, left_h = Inches(0.25), Inches(0.55), Inches(3.45), Inches(5.45)
    center_x, center_y, center_w, center_h = Inches(4.0), Inches(0.55), Inches(5.2), Inches(5.45)
    right_x, right_y, right_w, right_h = Inches(9.55), Inches(0.55), Inches(3.25), Inches(5.45)

    _panel(slide, left_x, left_y, left_w, left_h, "金融量化组合优化")
    _panel(slide, center_x, center_y, center_w, center_h, "量电混合量化投资平台")
    _panel(slide, right_x, right_y, right_w, right_h, "关键结果与应用价值")

    _textbox(
        slide,
        left_x + Inches(0.2),
        left_y + Inches(0.62),
        left_w - Inches(0.4),
        Inches(0.7),
        "面向多标的离散选股与资产配置，在收益、风险、行业暴露与持仓数量等约束下进行联合优化。",
        size=10,
        color=MUTED,
    )
    alert = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x + Inches(0.18), left_y + Inches(1.28), left_w - Inches(0.36), Inches(1.55))
    alert.fill.solid()
    alert.fill.fore_color.rgb = RGBColor(0x14, 0x7E, 0xC1)
    alert.line.color.rgb = RGBColor(0x59, 0xD1, 0xFF)
    _textbox(
        slide,
        left_x + Inches(0.25),
        left_y + Inches(1.38),
        left_w - Inches(0.5),
        Inches(0.4),
        "核心挑战：组合空间爆炸、全局最优难求、投研闭环迟滞",
        size=10,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    left_bullets = [
        "候选池稍有扩大，可行解规模即指数膨胀。",
        "500 选 30 的理论组合数已超过 10^48。",
        "传统启发式在高维协方差下易陷入局部最优。",
    ]
    for idx, item in enumerate(left_bullets):
        _textbox(
            slide,
            left_x + Inches(0.28),
            left_y + Inches(1.85 + 0.35 * idx),
            left_w - Inches(0.56),
            Inches(0.3),
            f"• {item}",
            size=9,
            color=TEXT,
        )

    chip_y = left_y + Inches(3.15)
    chip_w = Inches(0.88)
    for idx, (label, fill) in enumerate([
        ("50选10\n103亿", RGBColor(0xEB, 0xF1, 0xFF)),
        ("100选15\n10^17", RGBColor(0xD9, 0xE4, 0xFF)),
        ("500选30\n10^48", RGBColor(0xC8, 0xD7, 0xFF)),
    ]):
        _chip(slide, left_x + Inches(0.28 + idx * 0.96), chip_y, chip_w, Inches(0.75), label, fill=fill, size=10, bold=True)

    _textbox(
        slide,
        left_x + Inches(0.2),
        left_y + Inches(4.1),
        left_w - Inches(0.4),
        Inches(1.2),
        "量子赋能 + QUBO 建模 + 量电混合求解：\n将组合优化从“经验近似”推进到“离散优化可编排、可评估、可迭代”的新范式。",
        size=10,
        bold=True,
        color=RGBColor(0x70, 0xE8, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    # Center process map
    top_y = center_y + Inches(0.58)
    _chip(slide, center_x + Inches(0.2), top_y, Inches(2.15), Inches(0.45), "收益-风险-约束联合建模", fill=BLUE, text_color=TEXT, size=10, bold=True)
    _chip(slide, center_x + Inches(2.5), top_y, Inches(2.15), Inches(0.45), "QUBO 映射与惩罚项自适应", fill=RGBColor(0xC9, 0xD6, 0xFF), size=10, bold=True)

    mid_y = center_y + Inches(1.26)
    mid_labels = ["因子/行情输入", "协方差估计", "持仓约束", "工作流编排"]
    for idx, label in enumerate(mid_labels):
        _chip(slide, center_x + Inches(0.22 + idx * 1.18), mid_y, Inches(1.05), Inches(0.42), label, fill=RGBColor(0xEE, 0xF2, 0xFF), size=9)

    _chip(slide, center_x + Inches(0.25), center_y + Inches(2.0), Inches(2.05), Inches(0.38), "基于 QUBO 的组合优化目标构建", fill=RGBColor(0x6B, 0xA8, 0xFF), text_color=TEXT, size=9, bold=True)
    _chip(slide, center_x + Inches(2.55), center_y + Inches(2.0), Inches(2.05), Inches(0.38), "基于量电混合的并行求解与评估", fill=RGBColor(0xB9, 0xCC, 0xFF), size=9, bold=True)

    solve_y = center_y + Inches(2.62)
    solver_labels = ["经典基线", "模拟退火", "QAOA", "统一评估"]
    solver_fills = [RGBColor(0xF0, 0xF4, 0xFF), RGBColor(0xD9, 0xE5, 0xFF), BLUE, RGBColor(0xF0, 0xF4, 0xFF)]
    solver_text = [RGBColor(0x12, 0x1C, 0x30), RGBColor(0x12, 0x1C, 0x30), TEXT, RGBColor(0x12, 0x1C, 0x30)]
    for idx, label in enumerate(solver_labels):
        _chip(slide, center_x + Inches(0.28 + idx * 1.18), solve_y, Inches(1.02), Inches(0.42), label, fill=solver_fills[idx], text_color=solver_text[idx], size=9, bold=(label == "QAOA"))

    flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_x + Inches(0.22), center_y + Inches(3.18), center_w - Inches(0.44), Inches(0.9))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0x09, 0x0F, 0x24)
    flow_box.line.color.rgb = PANEL_BORDER
    flow_box.line.width = Pt(1.0)
    _textbox(
        slide,
        center_x + Inches(0.25),
        center_y + Inches(3.28),
        center_w - Inches(0.5),
        Inches(0.4),
        "数据接入 → 统计估计 → QUBO 构造 → QAOA / 退火求解 → 统一评估",
        size=10,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        center_x + Inches(0.25),
        center_y + Inches(3.70),
        center_w - Inches(0.5),
        Inches(0.3),
        "经典侧承载数据与回测，量子侧聚焦离散组合优化，形成可复用投研闭环。",
        size=9,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    # Right panel KPIs
    _textbox(slide, right_x + Inches(0.1), right_y + Inches(0.62), right_w - Inches(0.2), Inches(0.24), "固定种子仿真结果（6标的 / K=3 / 120日）", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    stat_titles = [("2.6%", "相对贪心夏普提升"), ("10^48", "500选30组合空间"), ("2~110ms", "多方法同口径求解")]
    stat_xs = [right_x + Inches(0.0), right_x + Inches(1.1), right_x + Inches(2.2)]
    for x, (big, small) in zip(stat_xs, stat_titles):
        _textbox(slide, x, right_y + Inches(1.0), Inches(1.05), Inches(0.42), big, size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        _textbox(slide, x, right_y + Inches(1.5), Inches(1.05), Inches(0.3), small, size=8, color=TEXT, align=PP_ALIGN.CENTER)

    _textbox(
        slide,
        right_x + Inches(0.2),
        right_y + Inches(1.82),
        right_w - Inches(0.4),
        Inches(0.54),
        "QAOA 夏普 1.20，较贪心 1.17 略优；在风险调整后收益维度展现更优解质量。",
        size=9,
        color=TEXT,
    )
    _textbox(
        slide,
        right_x + Inches(0.2),
        right_y + Inches(2.2),
        right_w - Inches(0.4),
        Inches(0.5),
        "通过量电混合求解与统一评估链路，支撑策略快速迭代与组合求解结果可解释输出。",
        size=8,
        color=MUTED,
    )
    _draw_chart(slide, right_x + Inches(0.25), right_y + Inches(2.82), Inches(2.72), Inches(1.55))
    _textbox(
        slide,
        right_x + Inches(0.2),
        right_y + Inches(4.55),
        right_w - Inches(0.4),
        Inches(0.52),
        "量子赋能 + 风险调整 + 组合优化：让量化投资从规则驱动走向更高维、更可编排的智能决策新范式。",
        size=10,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    _bottom_arrow(slide, Inches(0.45), Inches(6.12), Inches(3.2), Inches(0.44), "金融量化瓶颈：现象问题驱动研究")
    _bottom_arrow(slide, Inches(4.45), Inches(6.12), Inches(4.35), Inches(0.44), "量电混合量化投资平台 - 未来发展趋势")
    _bottom_arrow(slide, Inches(9.7), Inches(6.12), Inches(2.9), Inches(0.44), "量子赋能量化投资新范式")

    _textbox(
        slide,
        Inches(0.35),
        Inches(6.68),
        Inches(12.65),
        Inches(0.25),
        "说明：本页数据来自 FinQuanta 固定随机种子统计仿真，用于展示算法与工程链路，不构成投资收益承诺。",
        size=8,
        color=RGBColor(0x8C, 0x9B, 0xAF),
        align=PP_ALIGN.CENTER,
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print("Wrote", out)


if __name__ == "__main__":
    main()
